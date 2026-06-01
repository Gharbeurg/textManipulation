"""
correct_pptx.py

Objectif :
- Prendre un fichier PowerPoint .pptx en entrée
- Corriger les fautes d'orthographe du texte visible uniquement
- Utiliser un modèle local Ollama
- Générer un nouveau fichier PowerPoint corrigé
- Ne pas modifier volontairement la structure, les formes, les images ou la mise en page

Usage :
    python correct_pptx.py presentation.pptx

Sortie :
    presentation_corrige.pptx

Dépendances :
    pip install python-pptx requests

Prérequis :
    - Ollama doit être lancé
    - Le modèle défini dans OLLAMA_MODEL doit être installé
    - Vérifier le nom exact du modèle avec : ollama list
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Optional, Tuple

import requests
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


# ============================================================
# 1. CONFIGURATION
# ============================================================

# Le nom du modèle doit correspondre exactement au nom affiché par : ollama list
# Exemples possibles : "qwen3:14b", "qwen3:latest", "qwen2.5:14b", etc.
OLLAMA_MODEL = "qwen3.5"
CORRECTION_LANGUAGE = "fr"  # "fr" ou "en"
OLLAMA_TEMPERATURE = 0
OLLAMA_URL = "http://localhost:11434/api/chat"

# Filtrage prudent
# 2 mots permet de corriger des titres courts comme "Résultats principaux".
MIN_WORDS_TO_CORRECT = 2
MAX_UPPERCASE_RATIO = 0.80
MAX_SYMBOL_OR_DIGIT_RATIO = 0.70
MIN_LETTERS_TO_AVOID_SYMBOL_FILTER = 8

# Protection de la mise en forme
# Si True, le programme ignore les paragraphes composés de plusieurs runs.
# Un run est un morceau de texte avec sa propre mise en forme.
# Cela protège les paragraphes avec du gras, de l'italique ou des couleurs internes.
SKIP_MULTI_RUN_PARAGRAPHS = True

# Validation de la réponse du modèle
MAX_LENGTH_INCREASE_RATIO = 2.0
MIN_LENGTH_DECREASE_RATIO = 0.30
OLLAMA_TIMEOUT_SECONDS = 120

# Comportement console
PRINT_SKIPPED_TEXTS = True
PRINT_CORRECTIONS = True
PRINT_UNCHANGED_TEXTS = False


# ============================================================
# 2. COMPTEURS
# ============================================================

@dataclass
class ProcessingStats:
    slides_processed: int = 0
    text_blocks_seen: int = 0
    text_blocks_corrected: int = 0
    text_blocks_unchanged: int = 0
    text_blocks_skipped: int = 0
    errors: int = 0


# ============================================================
# 3. UTILITAIRES GÉNÉRAUX
# ============================================================

def build_output_path(input_path: Path) -> Path:
    """
    Génère automatiquement le chemin du fichier corrigé sans écraser
    un fichier déjà existant.

    Exemples :
        presentation.pptx -> presentation_corrige.pptx
        presentation.pptx -> presentation_corrige_2.pptx si le premier existe déjà
    """
    base_output_path = input_path.with_name(f"{input_path.stem}_corrige{input_path.suffix}")

    if not base_output_path.exists():
        return base_output_path

    counter = 2
    while True:
        candidate = input_path.with_name(f"{input_path.stem}_corrige_{counter}{input_path.suffix}")
        if not candidate.exists():
            return candidate
        counter += 1


def normalize_text_for_analysis(text: str) -> str:
    """
    Nettoie légèrement le texte pour décider s'il faut le corriger.

    Important :
    Cette fonction ne doit pas servir à remplacer le texte dans PowerPoint.
    Elle sert uniquement à l'analyse.
    """
    return text.strip()


def count_words(text: str) -> int:
    """
    Compte approximativement les mots.
    """
    words = re.findall(r"\b[\wÀ-ÿ'-]+\b", text, flags=re.UNICODE)
    return len(words)


def truncate_for_console(text: str, max_length: int = 300) -> str:
    """
    Évite d'afficher des blocs énormes dans la console.
    """
    clean = text.replace("\n", "\\n")
    if len(clean) <= max_length:
        return clean
    return clean[: max_length - 3] + "..."


def print_correction(slide_number: int, original: str, corrected: str) -> None:
    """
    Affiche une correction dans la console.
    """
    if not PRINT_CORRECTIONS:
        return

    print("-" * 80)
    print(f"Slide {slide_number}")
    print(f"Original : {truncate_for_console(original)}")
    print(f"Corrigé  : {truncate_for_console(corrected)}")


def print_skipped(slide_number: int, text: str, reason: str) -> None:
    """
    Affiche un texte ignoré dans la console.
    """
    if not PRINT_SKIPPED_TEXTS:
        return

    print("-" * 80)
    print(f"Slide {slide_number}")
    print(f"Ignoré : {truncate_for_console(text)}")
    print(f"Raison : {reason}")


def print_unchanged(slide_number: int, text: str) -> None:
    """
    Affiche un texte envoyé au modèle mais revenu inchangé.
    Désactivé par défaut pour garder la console lisible.
    """
    if not PRINT_UNCHANGED_TEXTS:
        return

    print("-" * 80)
    print(f"Slide {slide_number}")
    print(f"Inchangé : {truncate_for_console(text)}")


def print_error(slide_number: int, text: str, error: Exception) -> None:
    """
    Affiche une erreur de correction dans la console.
    """
    print("-" * 80)
    print(f"Slide {slide_number}")
    print(f"Erreur sur le texte : {truncate_for_console(text)}")
    print("Action : texte original conservé")
    print(f"Détail : {error}")


def print_summary(stats: ProcessingStats, output_path: Path) -> None:
    """
    Affiche le résumé final du traitement.
    """
    print("=" * 80)
    print("Traitement terminé")
    print(f"Slides parcourues        : {stats.slides_processed}")
    print(f"Blocs de texte vus       : {stats.text_blocks_seen}")
    print(f"Blocs corrigés           : {stats.text_blocks_corrected}")
    print(f"Blocs inchangés          : {stats.text_blocks_unchanged}")
    print(f"Blocs ignorés            : {stats.text_blocks_skipped}")
    print(f"Erreurs                  : {stats.errors}")
    print(f"Fichier généré           : {output_path}")
    print("=" * 80)


# ============================================================
# 4. FILTRAGE DES TEXTES À IGNORER
# ============================================================

def should_skip_text(text: str) -> Tuple[bool, Optional[str]]:
    """
    Décide si un texte doit être ignoré.

    Retourne :
        (True, raison) si le texte doit être ignoré
        (False, None) si le texte peut être envoyé au modèle
    """
    cleaned = normalize_text_for_analysis(text)

    if is_empty_or_whitespace(cleaned):
        return True, "texte vide"

    if is_too_short(cleaned):
        return True, "texte trop court ou ambigu"

    if looks_like_dose_or_reference(cleaned):
        return True, "dose, valeur scientifique ou référence"

    if is_mostly_uppercase(cleaned):
        return True, "texte majoritairement en majuscules"

    if is_mostly_numeric_or_symbols(cleaned):
        return True, "texte composé surtout de chiffres ou symboles"

    if looks_like_url_or_email(cleaned):
        return True, "URL ou adresse email"

    return False, None


def is_empty_or_whitespace(text: str) -> bool:
    """
    Ignore les textes vides.
    """
    return not text.strip()


def is_too_short(text: str) -> bool:
    """
    Ignore les textes très courts, souvent risqués à corriger.

    Exemples :
        HTA
        RGO
        IA
        10 mg

    Remarque :
    On garde un seuil bas pour ne pas ignorer les titres courts.
    """
    return count_words(text) < MIN_WORDS_TO_CORRECT


def is_mostly_uppercase(text: str) -> bool:
    """
    Ignore les textes majoritairement en majuscules.

    Exemples :
        MARKET ACCESS
        DIGITAL HEALTHCARE
        RGO / CONSTIPATION / PROBIOTIQUES

    On évite d'appliquer ce filtre sur des textes très courts,
    car quelques acronymes dans un titre peuvent fausser le ratio.
    """
    letters = [char for char in text if char.isalpha()]
    if len(letters) < 5:
        return False

    uppercase_letters = [char for char in letters if char.isupper()]
    ratio = len(uppercase_letters) / len(letters)
    return ratio >= MAX_UPPERCASE_RATIO


def is_mostly_numeric_or_symbols(text: str) -> bool:
    """
    Ignore les textes composés surtout de chiffres ou symboles.

    Exemples :
        2026
        12 %
        N = 124
        p < 0,05

    Le filtre reste prudent : si le texte contient assez de lettres,
    on ne l'exclut pas uniquement à cause de la ponctuation ou des espaces.
    """
    if not text:
        return True

    letters = [char for char in text if char.isalpha()]
    if len(letters) >= MIN_LETTERS_TO_AVOID_SYMBOL_FILTER:
        return False

    numeric_or_symbol_chars = [
        char for char in text
        if char.isdigit() or not char.isalpha()
    ]
    ratio = len(numeric_or_symbol_chars) / len(text)
    return ratio >= MAX_SYMBOL_OR_DIGIT_RATIO


def looks_like_dose_or_reference(text: str) -> bool:
    """
    Ignore les doses, références, valeurs scientifiques ou notations statistiques.
    """
    patterns = [
        r"\b\d+\s?(mg|g|kg|ml|l|µg|mcg|mmol|mol|ui|iu)\b",
        r"\b\d+\s?(fois|x)\s?/\s?(jour|j|day|d)\b",
        r"\bp\s?[<=>]\s?0[\.,]\d+\b",
        r"\bN\s?=\s?\d+\b",
        r"\bFigure\s+\d+\b",
        r"\bFig\.\s?\d+\b",
        r"\bTable\s+\d+\b",
        r"\bSlide\s+\d+\b",
        r"\bQ[1-4]\s?20\d{2}\b",
    ]

    return any(re.search(pattern, text, flags=re.IGNORECASE) for pattern in patterns)


def looks_like_url_or_email(text: str) -> bool:
    """
    Ignore les URL et adresses email.
    """
    patterns = [
        r"https?://",
        r"www\.",
        r"\b[\w\.-]+@[\w\.-]+\.\w+\b",
    ]
    return any(re.search(pattern, text, flags=re.IGNORECASE) for pattern in patterns)


# ============================================================
# 5. PROMPT ET APPEL OLLAMA
# ============================================================

def build_correction_prompt(text: str) -> str:
    """
    Construit le prompt envoyé au modèle.

    Le modèle doit répondre en JSON strict pour limiter les réponses bavardes.
    """
    if CORRECTION_LANGUAGE == "fr":
        instruction = """
Tu es un correcteur orthographique professionnel.
Corrige uniquement les fautes d'orthographe, de grammaire, d'accord et de ponctuation évidente.
Ne reformule pas.
Ne change pas le sens.
Ne raccourcis pas le texte.
N'allonge pas volontairement le texte.
Ne modifie pas les chiffres.
Ne modifie pas les acronymes.
Ne modifie pas les noms propres.
Ne modifie pas les termes métier.
Conserve la casse générale du texte.
Conserve les retours à la ligne s'il y en a.
Réponds uniquement avec un JSON valide, sans texte avant ni après.
""".strip()
    elif CORRECTION_LANGUAGE == "en":
        instruction = """
You are a professional proofreader.
Correct only spelling, grammar, agreement, and obvious punctuation mistakes.
Do not rephrase.
Do not change the meaning.
Do not shorten the text.
Do not intentionally expand the text.
Do not modify numbers.
Do not modify acronyms.
Do not modify proper nouns.
Do not modify business or medical terms.
Preserve the general casing of the text.
Preserve line breaks if any.
Reply only with valid JSON, with no text before or after.
""".strip()
    else:
        raise ValueError(f"Langue non supportée : {CORRECTION_LANGUAGE}")

    return f"""
{instruction}

Format attendu :
{{
  "corrected_text": "texte corrigé",
  "changed": true
}}

Texte à corriger :
{text}
""".strip()


def correct_text_with_ollama(text: str) -> str:
    """
    Envoie un texte à Ollama et récupère le texte corrigé.

    En cas de réponse invalide, une exception est levée.
    Le niveau supérieur décidera de conserver le texte original.
    """
    prompt = build_correction_prompt(text)

    payload = {
        "model": OLLAMA_MODEL,
        "messages": [
            {
                "role": "user",
                "content": prompt,
            }
        ],
        "stream": False,
        "format": "json",
        "options": {
            "temperature": OLLAMA_TEMPERATURE,
        },
    }

    response = requests.post(OLLAMA_URL, json=payload, timeout=OLLAMA_TIMEOUT_SECONDS)
    response.raise_for_status()

    data = response.json()
    raw_content = extract_ollama_message_content(data)

    corrected_text = parse_ollama_response(raw_content)
    corrected_text = normalize_corrected_text(original=text, corrected=corrected_text)
    validate_corrected_text(original=text, corrected=corrected_text)

    return corrected_text


def extract_ollama_message_content(data: dict[str, Any]) -> str:
    """
    Extrait le contenu texte depuis la réponse Ollama.
    """
    try:
        content = data["message"]["content"]
    except KeyError as error:
        raise ValueError(f"Réponse Ollama inattendue : champ manquant {error}") from error

    if not isinstance(content, str):
        raise ValueError("Réponse Ollama inattendue : le contenu n'est pas une chaîne")

    return content


def parse_ollama_response(raw_content: str) -> str:
    """
    Extrait corrected_text depuis la réponse JSON du modèle.

    Le modèle devrait répondre uniquement en JSON, mais un modèle local peut parfois
    ajouter du texte avant/après. Cette fonction tente donc d'extraire le premier
    bloc JSON trouvé dans la réponse.
    """
    content = raw_content.strip()

    # Nettoyage des blocs Markdown fréquents.
    content = content.removeprefix("```json").removeprefix("```").removesuffix("```").strip()

    # Si le modèle a ajouté du texte autour du JSON, on extrait le premier objet JSON.
    if not content.startswith("{"):
        match = re.search(r"\{.*\}", content, flags=re.DOTALL)
        if not match:
            raise ValueError("Réponse Ollama invalide : aucun JSON trouvé")
        content = match.group(0).strip()

    parsed = json.loads(content)

    if "corrected_text" not in parsed:
        raise ValueError("Réponse Ollama invalide : champ 'corrected_text' absent")

    corrected_text = parsed["corrected_text"]

    if not isinstance(corrected_text, str):
        raise ValueError("Réponse Ollama invalide : 'corrected_text' n'est pas une chaîne")

    return corrected_text


def normalize_corrected_text(original: str, corrected: str) -> str:
    """
    Nettoie prudemment la correction sans modifier le sens.

    On retire seulement les espaces ajoutés en début/fin, sauf si l'original
    contenait lui-même ces espaces. Comme PowerPoint stocke rarement des espaces
    significatifs en début/fin de paragraphe, ce comportement est acceptable.
    """
    return corrected.strip()


def validate_corrected_text(original: str, corrected: str) -> None:
    """
    Vérifie que la correction semble raisonnable.

    Si la correction semble dangereuse, on lève une exception.
    Le texte original sera alors conservé par le niveau supérieur.
    """
    original_stripped = original.strip()
    corrected_stripped = corrected.strip()

    if not corrected_stripped:
        raise ValueError("Correction vide")

    if len(original_stripped) == 0:
        raise ValueError("Texte original vide")

    length_ratio = len(corrected_stripped) / len(original_stripped)

    if length_ratio > MAX_LENGTH_INCREASE_RATIO:
        raise ValueError("Correction anormalement plus longue que l'original")

    if length_ratio < MIN_LENGTH_DECREASE_RATIO:
        raise ValueError("Correction anormalement plus courte que l'original")

    if added_wrapping_quotes(original_stripped, corrected_stripped):
        raise ValueError("La correction semble ajouter des guillemets enveloppants")

    if extract_numbers(original_stripped) != extract_numbers(corrected_stripped):
        raise ValueError("La correction modifie les chiffres")

    missing_acronyms = find_missing_acronyms(original_stripped, corrected_stripped)
    if missing_acronyms:
        joined_acronyms = ", ".join(missing_acronyms)
        raise ValueError(f"La correction supprime ou modifie des acronymes : {joined_acronyms}")

    suspicious_phrases = [
        "voici la correction",
        "here is the correction",
    ]

    lowered = corrected_stripped.lower()
    if any(phrase in lowered for phrase in suspicious_phrases):
        raise ValueError("La réponse semble contenir une explication au lieu du texte corrigé")


def added_wrapping_quotes(original: str, corrected: str) -> bool:
    """
    Détecte si le modèle a ajouté des guillemets autour du texte corrigé.

    Exemple à refuser :
        Original : Les patients sont suivis.
        Corrigé  : "Les patients sont suivis."
    """
    quote_pairs = [
        ('"', '"'),
        ("'", "'"),
        ("«", "»"),
        ("“", "”"),
    ]

    original_has_wrapping_quotes = any(
        original.startswith(open_quote) and original.endswith(close_quote)
        for open_quote, close_quote in quote_pairs
    )

    corrected_has_wrapping_quotes = any(
        corrected.startswith(open_quote) and corrected.endswith(close_quote)
        for open_quote, close_quote in quote_pairs
    )

    return corrected_has_wrapping_quotes and not original_has_wrapping_quotes


# ============================================================
# 6. TRAITEMENT POWERPOINT
# ============================================================

def process_presentation(input_path: Path, output_path: Path) -> ProcessingStats:
    """
    Ouvre le PowerPoint, parcourt les slides, corrige les textes visibles,
    puis sauvegarde le fichier corrigé.
    """
    stats = ProcessingStats()
    presentation = Presentation(str(input_path))

    for slide_index, slide in enumerate(presentation.slides, start=1):
        stats.slides_processed += 1
        process_slide(slide, slide_index, stats)

    presentation.save(str(output_path))
    return stats


def process_slide(slide: Any, slide_number: int, stats: ProcessingStats) -> None:
    """
    Parcourt toutes les formes d'une slide.
    """
    for shape in slide.shapes:
        process_shape(shape, slide_number, stats)


def process_shape(shape: Any, slide_number: int, stats: ProcessingStats) -> None:
    """
    Traite une forme PowerPoint.

    Cas gérés :
    - groupe de formes
    - forme avec texte
    - tableau
    - autres objets ignorés silencieusement
    """
    # 1. Groupes de formes : parcours récursif
    if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            process_shape(sub_shape, slide_number, stats)
        return

    # 2. Tableaux
    if safe_bool_attr(shape, "has_table"):
        process_table(shape.table, slide_number, stats)
        return

    # 3. Zones de texte / formes avec texte
    if safe_bool_attr(shape, "has_text_frame"):
        process_text_frame(shape.text_frame, slide_number, stats)
        return

    # 4. Images, graphiques, objets intégrés, etc. : ignorés
    return


def safe_bool_attr(obj: Any, attr_name: str) -> bool:
    """
    Lit prudemment un attribut booléen de python-pptx.
    """
    try:
        return bool(getattr(obj, attr_name, False))
    except Exception:
        return False


def process_table(table: Any, slide_number: int, stats: ProcessingStats) -> None:
    """
    Parcourt les cellules d'un tableau PowerPoint.
    """
    for row in table.rows:
        for cell in row.cells:
            if getattr(cell, "text_frame", None) is not None:
                process_text_frame(cell.text_frame, slide_number, stats)


def process_text_frame(text_frame: Any, slide_number: int, stats: ProcessingStats) -> None:
    """
    Corrige le texte visible dans une zone de texte.

    Le traitement se fait paragraphe par paragraphe.
    Dans PowerPoint, une puce est généralement un paragraphe.
    """
    for paragraph in text_frame.paragraphs:
        process_paragraph(paragraph, slide_number, stats)


def process_paragraph(paragraph: Any, slide_number: int, stats: ProcessingStats) -> None:
    """
    Corrige un paragraphe PowerPoint si le filtrage l'autorise.
    """
    original_text = get_paragraph_text(paragraph)
    stats.text_blocks_seen += 1

    if should_skip_paragraph_for_style(paragraph):
        stats.text_blocks_skipped += 1
        print_skipped(slide_number, original_text, "paragraphe avec plusieurs styles internes")
        return

    should_skip, reason = should_skip_text(original_text)
    if should_skip:
        stats.text_blocks_skipped += 1
        print_skipped(slide_number, original_text, reason or "raison inconnue")
        return

    try:
        corrected_text = correct_text_with_ollama(original_text)
    except Exception as error:
        stats.errors += 1
        print_error(slide_number, original_text, error)
        return

    if corrected_text == original_text:
        stats.text_blocks_unchanged += 1
        print_unchanged(slide_number, original_text)
        return

    replace_paragraph_text_simple(paragraph, corrected_text)
    stats.text_blocks_corrected += 1
    print_correction(slide_number, original_text, corrected_text)


def get_paragraph_text(paragraph: Any) -> str:
    """
    Récupère le texte d'un paragraphe PowerPoint.

    paragraph.text existe normalement dans python-pptx.
    On reconstruit depuis les runs si nécessaire.
    """
    text = getattr(paragraph, "text", None)
    if isinstance(text, str):
        return text

    runs = getattr(paragraph, "runs", [])
    return "".join(getattr(run, "text", "") for run in runs)


def should_skip_paragraph_for_style(paragraph: Any) -> bool:
    """
    Décide si un paragraphe doit être ignoré pour protéger sa mise en forme.

    Si SKIP_MULTI_RUN_PARAGRAPHS = True, on ignore les paragraphes qui ont
    plusieurs runs non vides. Cela évite d'écraser des styles internes.
    """
    if not SKIP_MULTI_RUN_PARAGRAPHS:
        return False

    runs = list(getattr(paragraph, "runs", []))
    non_empty_runs = [run for run in runs if getattr(run, "text", "")]

    return len(non_empty_runs) > 1


def replace_paragraph_text_simple(paragraph: Any, corrected_text: str) -> None:
    """
    Remplace le texte d'un paragraphe en conservant au mieux
    la mise en forme principale.

    Limite :
    Si le paragraphe contient plusieurs runs avec des styles différents,
    les différences internes peuvent être perdues.
    """
    runs = list(getattr(paragraph, "runs", []))

    if not runs:
        # Cas rare : paragraphe sans run existant.
        # python-pptx permet généralement d'utiliser paragraph.text.
        paragraph.text = corrected_text
        return

    # On garde la mise en forme du premier run.
    runs[0].text = corrected_text

    # On vide les autres runs pour éviter de laisser l'ancien texte.
    for run in runs[1:]:
        run.text = ""


# ============================================================
# 7. PROGRAMME PRINCIPAL
# ============================================================

def parse_args() -> argparse.Namespace:
    """
    Lit les arguments de ligne de commande.
    """
    parser = argparse.ArgumentParser(
        description="Corrige automatiquement l'orthographe d'un fichier PowerPoint avec Ollama."
    )
    parser.add_argument(
        "input_pptx",
        help="Chemin du fichier PowerPoint .pptx à corriger",
    )
    return parser.parse_args()


def validate_input_file(input_path: Path) -> None:
    """
    Vérifie le fichier d'entrée.
    """
    if not input_path.exists():
        raise FileNotFoundError(f"Fichier introuvable : {input_path}")

    if not input_path.is_file():
        raise ValueError(f"Le chemin n'est pas un fichier : {input_path}")

    if input_path.suffix.lower() != ".pptx":
        raise ValueError("Le fichier d'entrée doit être un .pptx")


def print_start_banner(input_path: Path, output_path: Path) -> None:
    """
    Affiche les paramètres de lancement.
    """
    print("=" * 80)
    print("Correction PowerPoint")
    print(f"Fichier source : {input_path}")
    print(f"Fichier sortie : {output_path}")
    print(f"Modèle Ollama  : {OLLAMA_MODEL}")
    print(f"Langue         : {CORRECTION_LANGUAGE}")
    print("=" * 80)


def main() -> int:
    """
    Point principal du programme.
    """
    args = parse_args()
    input_path = Path(args.input_pptx).expanduser().resolve()

    try:
        validate_input_file(input_path)
    except Exception as error:
        print(f"Erreur : {error}", file=sys.stderr)
        return 1

    output_path = build_output_path(input_path)
    print_start_banner(input_path, output_path)

    try:
        stats = process_presentation(input_path, output_path)
    except requests.exceptions.ConnectionError:
        print(
            "Erreur critique : impossible de joindre Ollama. "
            "Vérifie qu'Ollama est lancé et accessible sur http://localhost:11434.",
            file=sys.stderr,
        )
        return 1
    except requests.exceptions.HTTPError as error:
        print(
            "Erreur critique : Ollama a renvoyé une erreur HTTP. "
            f"Détail : {error}",
            file=sys.stderr,
        )
        return 1
    except Exception as error:
        print(f"Erreur critique : {error}", file=sys.stderr)
        return 1

    print_summary(stats, output_path)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
