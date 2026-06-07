"""
Correcteur orthographique et grammatical de fichiers PowerPoint (.pptx).

VERSION 1 AMÉLIORÉE
- Parcourt uniquement le répertoire configuré.
- Détecte la langue à partir des 3 premières diapositives.
- Traite uniquement le français et l'anglais.
- Corrige les zones de texte, titres, tableaux et objets groupés.
- Préserve au mieux la mise en forme et les liens hypertextes.
- Ignore les suggestions de style.
- Travaille sur un fichier temporaire.
- Remplace l'original uniquement après validation complète.
- Ne réécrit pas l'original lorsqu'aucune correction n'est appliquée.
- Préserve les sauts de ligne manuels.
- Évite le double traitement des cellules fusionnées.

Non traité dans cette version :
- notes du présentateur ;
- SmartArt non accessible par python-pptx ;
- texte interne des graphiques ;
- commentaires ;
- masques et dispositions ;
- texte présent dans les images ;
- objets OLE et équations.
"""

from __future__ import annotations

import copy
import os
import re
import shutil
import sys
import tempfile
import zipfile
from dataclasses import dataclass, field
from difflib import SequenceMatcher
from pathlib import Path
from typing import Any, Iterable, Optional

import language_tool_python
from langdetect import DetectorFactory, LangDetectException, detect
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


# =============================================================================
# 1. CONFIGURATION À MODIFIER
# =============================================================================

# Répertoire contenant les fichiers .pptx à corriger.
POWERPOINT_DIRECTORY = Path(r"C:\PYTHON\.entree\SourcesPowerpoint")

# Fichier contenant un mot ou une expression protégée par ligne.
PROTECTED_WORDS_FILE = Path(r"C:\PYTHON\.params\mots_proteges.txt")

# Détection de langue.
LANGUAGE_DETECTION_SLIDES = 3
MIN_LANGUAGE_CHARACTERS = 300

# Langues autorisées.
SUPPORTED_LANGUAGES = {
    "fr": "fr-FR",
    "en": "en-US",
}

# Extensions et fichiers ignorés.
POWERPOINT_EXTENSION = ".pptx"
POWERPOINT_TEMPORARY_PREFIX = "~$"

# Les catégories ou règles contenant ces termes sont ignorées.
# Elles correspondent en général à des conseils de style ou de reformulation.
EXCLUDED_CATEGORY_KEYWORDS = {
    "STYLE",
    "REDUNDANCY",
    "PLAIN_ENGLISH",
    "SEMANTICS",
}

EXCLUDED_RULE_ID_KEYWORDS = {
    "STYLE",
    "REDUNDANCY",
    "SENTENCE_LENGTH",
    "PASSIVE_VOICE",
    "WORDINESS",
    "READABILITY",
}

EXCLUDED_MESSAGE_KEYWORDS_FR = {
    "style",
    "reformul",
    "phrase longue",
    "répétition",
    "redondan",
    "plus concis",
    "plus simple",
}

EXCLUDED_MESSAGE_KEYWORDS_EN = {
    "style",
    "rephrase",
    "long sentence",
    "repetition",
    "redundant",
    "more concise",
    "simpler wording",
}

# Types de problèmes LanguageTool considérés comme suffisamment sûrs.
ALLOWED_ISSUE_TYPES = {
    "misspelling",
    "grammar",
    "grammatical",
    "typographical",
    "duplication",
}

# Catégories acceptées lorsqu'elles sont disponibles.
ALLOWED_CATEGORY_KEYWORDS = {
    "GRAMMAR",
    "AGREEMENT",
    "TYPOS",
    "TYPOGRAPHY",
    "PUNCTUATION",
    "CASING",
    "SPELLING",
}

# Rend la détection de langue reproductible.
DetectorFactory.seed = 0


# =============================================================================
# 2. EXPRESSIONS RÉGULIÈRES
# =============================================================================

URL_PATTERN = re.compile(
    r"""(?ix)
    \b(
        https?://[^\s<>"']+
        |
        www\.[^\s<>"']+
    )
    """
)

EMAIL_PATTERN = re.compile(
    r"\b[A-Z0-9._%+-]+@[A-Z0-9.-]+\.[A-Z]{2,}\b",
    re.IGNORECASE,
)

WINDOWS_PATH_PATTERN = re.compile(
    r"""(?ix)
    (?:
        [A-Z]:\\
        |
        \\\\
    )
    [^\s<>"|?*]+
    """
)

UNIX_PATH_PATTERN = re.compile(
    r"(?<!\w)/(?:[^/\s]+/)*[^/\s]+"
)

XML_TAG_PATTERN = re.compile(r"<[^>\r\n]+>")

# Codes contenant au moins une lettre et un chiffre : ABC123, H1N1, V2.0...
ALPHANUMERIC_CODE_PATTERN = re.compile(
    r"\b(?=\w*[A-Za-z])(?=\w*\d)[A-Za-z0-9._/-]+\b"
)

# Valeurs numériques, avec décimales, pourcentages et unités simples.
NUMBER_PATTERN = re.compile(
    r"""(?ix)
    (?<!\w)
    [+-]?
    \d+(?:[.,]\d+)?
    (?:\s?[%°])?
    (?:\s?(?:mg|g|kg|ml|l|cm|mm|m|km|hz|khz|mhz|ghz|€|\$))?
    (?!\w)
    """
)

# Formules simples ou fragments de code.
FORMULA_PATTERN = re.compile(
    r"""(?x)
    (?:
        \b[A-Za-z_][A-Za-z0-9_]*\s*(?:==|!=|<=|>=|=|\+|-|\*|/)\s*
        [A-Za-z0-9_.()+\-*/]+\b
    )
    """
)

LETTER_PATTERN = re.compile(r"[A-Za-zÀ-ÖØ-öø-ÿ]")

# python-pptx représente généralement un saut de ligne manuel par \x0b.
SOFT_LINE_BREAK = "\x0b"


# =============================================================================
# 3. STRUCTURES DE DONNÉES
# =============================================================================

@dataclass
class FileStatistics:
    filename: str
    detected_language: Optional[str] = None
    slides_processed: int = 0
    corrections_applied: int = 0
    corrections_ignored: int = 0
    protected_corrections: int = 0
    style_corrections_ignored: int = 0
    corrections_without_suggestion: int = 0
    errors: int = 0
    status: str = "PENDING"
    skip_reason: str = ""


@dataclass
class TextCorrectionResult:
    original_text: str
    corrected_text: str
    applied_count: int = 0
    ignored_count: int = 0
    protected_count: int = 0
    style_ignored_count: int = 0
    without_suggestion_count: int = 0

    @property
    def changed(self) -> bool:
        return self.original_text != self.corrected_text


@dataclass
class RunStyle:
    """
    Copie des propriétés XML d'un run.

    La copie du nœud <a:rPr> permet de conserver plus fidèlement :
    - police ;
    - taille ;
    - couleur ;
    - gras, italique et soulignement ;
    - langue ;
    - exposant ou indice ;
    - lien hypertexte.
    """
    rpr_xml: Optional[Any] = None


@dataclass
class ParagraphSnapshot:
    """
    Propriétés du paragraphe à restaurer si nécessaire.
    paragraph.clear() conserve normalement <a:pPr>, mais une copie est gardée
    par sécurité.
    """
    ppr_xml: Optional[Any] = None


# =============================================================================
# 4. OUTILS GÉNÉRAUX
# =============================================================================

def normalize_text(value: Any) -> str:
    if value is None:
        return ""
    return str(value)


def safe_getattr(obj: Any, attribute: str, default: Any = None) -> Any:
    try:
        return getattr(obj, attribute, default)
    except Exception:
        return default


def merge_ranges(ranges: Iterable[tuple[int, int]]) -> list[tuple[int, int]]:
    cleaned = sorted(
        (max(0, start), max(0, end))
        for start, end in ranges
        if end > start
    )

    if not cleaned:
        return []

    merged: list[tuple[int, int]] = [cleaned[0]]

    for start, end in cleaned[1:]:
        previous_start, previous_end = merged[-1]

        if start <= previous_end:
            merged[-1] = (previous_start, max(previous_end, end))
        else:
            merged.append((start, end))

    return merged


def ranges_overlap(
    start: int,
    end: int,
    protected_ranges: Iterable[tuple[int, int]],
) -> bool:
    return any(start < protected_end and end > protected_start
               for protected_start, protected_end in protected_ranges)


# =============================================================================
# 5. VALIDATION DE LA CONFIGURATION
# =============================================================================

def validate_configuration() -> None:
    errors: list[str] = []

    if not POWERPOINT_DIRECTORY.exists():
        errors.append(
            f"Le répertoire PowerPoint n'existe pas :\n{POWERPOINT_DIRECTORY}"
        )
    elif not POWERPOINT_DIRECTORY.is_dir():
        errors.append(
            f"Le chemin PowerPoint n'est pas un répertoire :\n"
            f"{POWERPOINT_DIRECTORY}"
        )

    if not PROTECTED_WORDS_FILE.exists():
        errors.append(
            f"Le fichier de mots protégés est introuvable :\n"
            f"{PROTECTED_WORDS_FILE}"
        )
    elif not PROTECTED_WORDS_FILE.is_file():
        errors.append(
            f"Le chemin des mots protégés n'est pas un fichier :\n"
            f"{PROTECTED_WORDS_FILE}"
        )

    if errors:
        raise RuntimeError("\n\n".join(errors))


# =============================================================================
# 6. MOTS ET EXPRESSIONS PROTÉGÉS
# =============================================================================

def load_protected_terms() -> list[str]:
    try:
        lines = PROTECTED_WORDS_FILE.read_text(encoding="utf-8-sig").splitlines()
    except OSError as exc:
        raise RuntimeError(
            f"Impossible de lire le fichier de mots protégés : "
            f"{PROTECTED_WORDS_FILE}"
        ) from exc

    terms: list[str] = []
    seen: set[str] = set()

    for line in lines:
        value = line.strip()

        if not value or value.startswith("#"):
            continue

        key = value.casefold()

        if key not in seen:
            terms.append(value)
            seen.add(key)

    return terms


def protected_term_pattern(term: str) -> re.Pattern[str]:
    escaped = re.escape(term)

    # Les limites basées sur \w empêchent "patient" de protéger "patients".
    return re.compile(
        rf"(?<!\w){escaped}(?!\w)",
        re.IGNORECASE,
    )


def find_technical_ranges(text: str) -> list[tuple[int, int]]:
    ranges: list[tuple[int, int]] = []

    patterns = (
        URL_PATTERN,
        EMAIL_PATTERN,
        WINDOWS_PATH_PATTERN,
        UNIX_PATH_PATTERN,
        XML_TAG_PATTERN,
        ALPHANUMERIC_CODE_PATTERN,
        NUMBER_PATTERN,
        FORMULA_PATTERN,
    )

    for pattern in patterns:
        ranges.extend((match.start(), match.end())
                      for match in pattern.finditer(text))

    return merge_ranges(ranges)


def find_protected_ranges(
    text: str,
    protected_terms: list[str],
) -> list[tuple[int, int]]:
    ranges = find_technical_ranges(text)

    for term in protected_terms:
        pattern = protected_term_pattern(term)
        ranges.extend((match.start(), match.end())
                      for match in pattern.finditer(text))

    return merge_ranges(ranges)


# =============================================================================
# 7. RECHERCHE ET CONTRÔLE DES FICHIERS
# =============================================================================

def find_powerpoint_files() -> list[Path]:
    files = [
        path
        for path in POWERPOINT_DIRECTORY.iterdir()
        if (
            path.is_file()
            and path.suffix.lower() == POWERPOINT_EXTENSION
            and not path.name.startswith(POWERPOINT_TEMPORARY_PREFIX)
        )
    ]

    return sorted(files, key=lambda path: path.name.casefold())


def has_powerpoint_lock_file(file_path: Path) -> bool:
    lock_file = file_path.with_name(f"~${file_path.name}")
    return lock_file.exists()


def can_open_for_exclusive_write(file_path: Path) -> bool:
    """
    Tente un verrouillage non bloquant sous Windows.

    Sur les autres systèmes, l'ouverture en lecture/écriture est utilisée
    comme contrôle minimal.
    """
    try:
        with file_path.open("r+b") as stream:
            if os.name != "nt":
                return True

            import msvcrt

            stream.seek(0)
            try:
                msvcrt.locking(stream.fileno(), msvcrt.LK_NBLCK, 1)
                msvcrt.locking(stream.fileno(), msvcrt.LK_UNLCK, 1)
                return True
            except OSError:
                return False

    except OSError:
        return False


def is_valid_pptx_archive(file_path: Path) -> tuple[bool, str]:
    required_files = {
        "[Content_Types].xml",
        "_rels/.rels",
        "ppt/presentation.xml",
        "ppt/_rels/presentation.xml.rels",
    }

    try:
        with zipfile.ZipFile(file_path, "r") as archive:
            bad_file = archive.testzip()

            if bad_file is not None:
                return False, f"archive ZIP endommagée : {bad_file}"

            names = set(archive.namelist())
            missing = sorted(required_files - names)

            if missing:
                return False, (
                    "composants PowerPoint manquants : "
                    + ", ".join(missing)
                )

            if not any(
                name.startswith("ppt/slides/slide")
                and name.endswith(".xml")
                for name in names
            ):
                return False, "aucune diapositive trouvée"

    except zipfile.BadZipFile:
        return False, "le fichier n'est pas une archive PowerPoint valide"
    except OSError as exc:
        return False, f"lecture impossible : {exc}"

    return True, ""


def can_process_file(file_path: Path) -> tuple[bool, str]:
    if has_powerpoint_lock_file(file_path):
        return False, "fichier ouvert ou verrouillé"

    if not can_open_for_exclusive_write(file_path):
        return False, "fichier ouvert, verrouillé ou non modifiable"

    is_valid, reason = is_valid_pptx_archive(file_path)

    if not is_valid:
        return False, reason

    try:
        Presentation(str(file_path))
    except Exception as exc:
        return False, (
            "présentation protégée, chiffrée ou illisible "
            f"({type(exc).__name__})"
        )

    return True, ""


# =============================================================================
# 8. EXTRACTION DU TEXTE ET DÉTECTION DE LANGUE
# =============================================================================

def iter_shape_texts(shape: Any) -> Iterable[str]:
    """
    Extrait le texte accessible avec python-pptx sans modifier la présentation.
    """
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child_shape in shape.shapes:
            yield from iter_shape_texts(child_shape)
        return

    if safe_getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    if paragraph.text:
                        yield paragraph.text
        return

    if safe_getattr(shape, "has_text_frame", False):
        for paragraph in shape.text_frame.paragraphs:
            if paragraph.text:
                yield paragraph.text


def extract_language_detection_text(presentation: Presentation) -> str:
    texts: list[str] = []

    maximum = min(LANGUAGE_DETECTION_SLIDES, len(presentation.slides))

    for slide_index in range(maximum):
        slide = presentation.slides[slide_index]

        for shape in slide.shapes:
            texts.extend(iter_shape_texts(shape))

    return "\n".join(texts)


def clean_text_for_language_detection(text: str) -> str:
    characters = list(text)

    # Remplace les plages techniques par des espaces afin de ne pas coller
    # artificiellement les mots situés de part et d'autre.
    for start, end in find_technical_ranges(text):
        for index in range(start, min(end, len(characters))):
            characters[index] = " "

    cleaned = "".join(characters)
    cleaned = re.sub(r"\s+", " ", cleaned).strip()
    return cleaned


def detect_presentation_language(text: str) -> tuple[Optional[str], str]:
    cleaned = clean_text_for_language_detection(text)

    if len(cleaned) < MIN_LANGUAGE_CHARACTERS:
        return None, (
            f"moins de {MIN_LANGUAGE_CHARACTERS} caractères utiles dans les "
            f"{LANGUAGE_DETECTION_SLIDES} premières diapositives"
        )

    try:
        detected_language = detect(cleaned)
    except LangDetectException:
        return None, "langue non détectable"

    language_tool_code = SUPPORTED_LANGUAGES.get(detected_language)

    if language_tool_code is None:
        return None, (
            f"langue non prise en charge : {detected_language}"
        )

    return language_tool_code, ""


# =============================================================================
# 9. LANGUAGETOOL
# =============================================================================

def initialize_language_tools() -> dict[str, Any]:
    tools: dict[str, Any] = {}

    try:
        for language_code in SUPPORTED_LANGUAGES.values():
            tools[language_code] = language_tool_python.LanguageTool(
                language_code
            )
    except Exception:
        close_language_tools(tools)
        raise

    return tools


def close_language_tools(language_tools: dict[str, Any]) -> None:
    for tool in language_tools.values():
        try:
            tool.close()
        except Exception:
            pass


def match_category(match: Any) -> str:
    category = safe_getattr(match, "category", "")

    if hasattr(category, "id"):
        category = safe_getattr(category, "id", category)

    return normalize_text(category).upper()


def match_rule_id(match: Any) -> str:
    return normalize_text(
        safe_getattr(
            match,
            "rule_id",
            safe_getattr(match, "ruleId", ""),
        )
    ).upper()


def match_issue_type(match: Any) -> str:
    return normalize_text(
        safe_getattr(
            match,
            "rule_issue_type",
            safe_getattr(match, "ruleIssueType", ""),
        )
    ).lower()


def match_replacements(match: Any) -> list[str]:
    replacements = safe_getattr(match, "replacements", [])

    if replacements is None:
        return []

    return [normalize_text(item) for item in replacements]


def match_offset(match: Any) -> int:
    return int(safe_getattr(match, "offset", 0))


def match_error_length(match: Any) -> int:
    return int(
        safe_getattr(
            match,
            "error_length",
            safe_getattr(match, "errorLength", 0),
        )
    )


def is_style_match(match: Any) -> bool:
    category = match_category(match)
    rule_id = match_rule_id(match)
    message = normalize_text(safe_getattr(match, "message", "")).casefold()

    if any(keyword in category for keyword in EXCLUDED_CATEGORY_KEYWORDS):
        return True

    if any(keyword in rule_id for keyword in EXCLUDED_RULE_ID_KEYWORDS):
        return True

    message_keywords = (
        EXCLUDED_MESSAGE_KEYWORDS_FR | EXCLUDED_MESSAGE_KEYWORDS_EN
    )

    return any(keyword.casefold() in message for keyword in message_keywords)


def is_allowed_language_tool_match(match: Any) -> bool:
    if is_style_match(match):
        return False

    issue_type = match_issue_type(match)
    category = match_category(match)
    rule_id = match_rule_id(match)

    if issue_type in ALLOWED_ISSUE_TYPES:
        return True

    # Accepte explicitement les règles grammaticales et d’accord.
    if any(
        keyword in category
        for keyword in {
            "GRAMMAR",
            "GRAMMAIRE",
            "AGREEMENT",
            "CONCORDANCE",
            "CONCORDANCES",
        }
    ):
        return True

    if any(
        keyword in rule_id
        for keyword in {
            "AGREEMENT",
            "CONCORD",
            "NOUN",
            "DET",
        }
    ):
        return True

    return any(
        keyword in category
        for keyword in ALLOWED_CATEGORY_KEYWORDS
    )


# =============================================================================
# 10. CORRECTION DU TEXTE
# =============================================================================
import unicodedata
from difflib import SequenceMatcher


def remove_accents(text: str) -> str:
    normalized = unicodedata.normalize("NFD", text)

    return "".join(
        character
        for character in normalized
        if unicodedata.category(character) != "Mn"
    )


def choose_best_replacement(
    original_fragment: str,
    replacements: list[str],
) -> str | None:
    if not replacements:
        return None

    original_without_accents = remove_accents(
        original_fragment
    ).casefold()

    def score(replacement: str) -> tuple[int, float, int]:
        replacement_without_accents = remove_accents(
            replacement
        ).casefold()

        # Priorité maximale lorsque la différence porte seulement
        # sur les accents ou les majuscules.
        accent_or_case_only = int(
            original_without_accents
            == replacement_without_accents
        )

        similarity = SequenceMatcher(
            None,
            original_fragment.casefold(),
            replacement.casefold(),
            autojunk=False,
        ).ratio()

        # À score égal, préférer la suggestion la plus courte.
        length_difference = -abs(
            len(original_fragment) - len(replacement)
        )

        return (
            accent_or_case_only,
            similarity,
            length_difference,
        )

    return max(replacements, key=score)

def _correct_text_segment(
    text: str,
    language_tool: Any,
    protected_terms: list[str],
) -> TextCorrectionResult:
    """
    Corrige un segment ne contenant pas de saut de ligne manuel.
    """
    result = TextCorrectionResult(
        original_text=text,
        corrected_text=text,
    )

    if not text.strip() or not LETTER_PATTERN.search(text):
        return result

    protected_ranges = find_protected_ranges(text, protected_terms)

    matches = language_tool.check(text)

    accepted: list[tuple[int, int, str]] = []

    for match in matches:
        start = match_offset(match)
        end = start + match_error_length(match)
        replacements = match_replacements(match)

        if is_style_match(match):
            result.style_ignored_count += 1
            result.ignored_count += 1
            continue

        if not is_allowed_language_tool_match(match):
            result.ignored_count += 1
            continue

        if not replacements:
            result.without_suggestion_count += 1
            result.ignored_count += 1
            continue

        if ranges_overlap(start, end, protected_ranges):
            result.protected_count += 1
            result.ignored_count += 1
            continue

        original_fragment = text[start:end]

        replacement = choose_best_replacement(
            original_fragment,
            replacements,
        )

        if replacement is None:
            result.without_suggestion_count += 1
            result.ignored_count += 1
            continue

        if replacement == text[start:end]:
            continue

        accepted.append((start, end, replacement))

    accepted.sort(key=lambda item: (item[0], item[1]), reverse=True)

    corrected = text
    occupied_ranges: list[tuple[int, int]] = []

    for start, end, replacement in accepted:
        if ranges_overlap(start, end, occupied_ranges):
            result.ignored_count += 1
            continue

        corrected = corrected[:start] + replacement + corrected[end:]
        occupied_ranges.append((start, end))
        result.applied_count += 1

    result.corrected_text = corrected
    return result


def correct_text(
    text: str,
    language_tool: Any,
    protected_terms: list[str],
) -> TextCorrectionResult:
    """
    Corrige le texte sans jamais supprimer ni déplacer les sauts de ligne
    manuels. Chaque ligne interne est contrôlée séparément, puis les séparateurs
    d'origine sont replacés à l'identique.
    """
    if SOFT_LINE_BREAK not in text and "\n" not in text:
        return _correct_text_segment(text, language_tool, protected_terms)

    separators = re.split(r"(\x0b|\n)", text)
    corrected_parts: list[str] = []

    final_result = TextCorrectionResult(
        original_text=text,
        corrected_text=text,
    )

    for part in separators:
        if part in {SOFT_LINE_BREAK, "\n"}:
            corrected_parts.append(part)
            continue

        segment_result = _correct_text_segment(
            part,
            language_tool,
            protected_terms,
        )

        corrected_parts.append(segment_result.corrected_text)
        final_result.applied_count += segment_result.applied_count
        final_result.ignored_count += segment_result.ignored_count
        final_result.protected_count += segment_result.protected_count
        final_result.style_ignored_count += segment_result.style_ignored_count
        final_result.without_suggestion_count += (
            segment_result.without_suggestion_count
        )

    final_result.corrected_text = "".join(corrected_parts)
    return final_result

# =============================================================================
# 11. CONSERVATION DE LA MISE EN FORME
# =============================================================================

def capture_paragraph_properties(paragraph: Any) -> ParagraphSnapshot:
    ppr = safe_getattr(paragraph._p, "pPr", None)

    return ParagraphSnapshot(
        ppr_xml=copy.deepcopy(ppr) if ppr is not None else None
    )


def restore_paragraph_properties(
    paragraph: Any,
    snapshot: ParagraphSnapshot,
) -> None:
    if snapshot.ppr_xml is None:
        return

    current_ppr = safe_getattr(paragraph._p, "pPr", None)

    if current_ppr is not None:
        paragraph._p.remove(current_ppr)

    paragraph._p.insert(0, copy.deepcopy(snapshot.ppr_xml))


def capture_run_style(run: Any) -> RunStyle:
    rpr = safe_getattr(run._r, "rPr", None)

    return RunStyle(
        rpr_xml=copy.deepcopy(rpr) if rpr is not None else None
    )


def build_character_style_map(
    paragraph: Any,
) -> tuple[str, list[RunStyle], list[int]]:
    """
    Construit une carte de styles à partir des runs tout en utilisant
    paragraph.text comme source de vérité, car celui-ci inclut les sauts de
    ligne manuels que paragraph.runs n'expose pas toujours.
    """
    paragraph_text = paragraph.text or ""
    runs = list(paragraph.runs)

    if not runs:
        return paragraph_text, [RunStyle()], [0] * len(paragraph_text)

    styles = [capture_run_style(run) for run in runs]
    run_texts = [run.text or "" for run in runs]

    character_style_ids: list[int] = []
    run_index = 0
    run_offset = 0
    last_style_id = 0

    for character in paragraph_text:
        if character in {SOFT_LINE_BREAK, "\n"}:
            character_style_ids.append(last_style_id)
            continue

        while run_index < len(run_texts):
            current_run_text = run_texts[run_index]

            if run_offset < len(current_run_text):
                break

            run_index += 1
            run_offset = 0

        if run_index >= len(run_texts):
            character_style_ids.append(last_style_id)
            continue

        character_style_ids.append(run_index)
        last_style_id = run_index
        run_offset += 1

    return paragraph_text, styles, character_style_ids

def dominant_style_id(
    style_ids: list[int],
    default_style_id: int = 0,
) -> int:
    if not style_ids:
        return default_style_id

    counts: dict[int, int] = {}

    for style_id in style_ids:
        counts[style_id] = counts.get(style_id, 0) + 1

    return max(counts, key=counts.get)


def build_corrected_style_map(
    original_text: str,
    corrected_text: str,
    original_style_ids: list[int],
) -> list[int]:
    if not corrected_text:
        return []

    if not original_style_ids:
        return [0] * len(corrected_text)

    matcher = SequenceMatcher(
        None,
        original_text,
        corrected_text,
        autojunk=False,
    )

    corrected_style_ids: list[Optional[int]] = [None] * len(corrected_text)

    for tag, i1, i2, j1, j2 in matcher.get_opcodes():
        if tag == "equal":
            for offset in range(j2 - j1):
                corrected_style_ids[j1 + offset] = original_style_ids[i1 + offset]

        elif tag == "replace":
            source_styles = original_style_ids[i1:i2]

            left_style = (
                original_style_ids[i1 - 1]
                if i1 > 0
                else original_style_ids[min(i1, len(original_style_ids) - 1)]
            )
            replacement_style = dominant_style_id(
                source_styles,
                default_style_id=left_style,
            )

            for index in range(j1, j2):
                corrected_style_ids[index] = replacement_style

        elif tag == "insert":
            if i1 > 0:
                insertion_style = original_style_ids[i1 - 1]
            elif i1 < len(original_style_ids):
                insertion_style = original_style_ids[i1]
            else:
                insertion_style = 0

            for index in range(j1, j2):
                corrected_style_ids[index] = insertion_style

        # "delete" ne crée aucun caractère dans le texte corrigé.

    # Remplit les rares valeurs non attribuées avec le style voisin.
    last_style = original_style_ids[0]

    for index, style_id in enumerate(corrected_style_ids):
        if style_id is None:
            corrected_style_ids[index] = last_style
        else:
            last_style = style_id

    return [int(style_id) for style_id in corrected_style_ids]


def apply_run_style(run: Any, style: RunStyle) -> None:
    if style.rpr_xml is None:
        return

    current_rpr = safe_getattr(run._r, "rPr", None)

    if current_rpr is not None:
        run._r.remove(current_rpr)

    run._r.insert(0, copy.deepcopy(style.rpr_xml))


def group_text_by_style(
    text: str,
    style_ids: list[int],
) -> list[tuple[str, int]]:
    if not text:
        return []

    if not style_ids:
        return [(text, 0)]

    groups: list[tuple[str, int]] = []
    start = 0
    current_style = style_ids[0]

    for index in range(1, len(text)):
        if style_ids[index] != current_style:
            groups.append((text[start:index], current_style))
            start = index
            current_style = style_ids[index]

    groups.append((text[start:], current_style))
    return groups


def rebuild_paragraph_with_styles(
    paragraph: Any,
    original_text: str,
    corrected_text: str,
    styles: list[RunStyle],
    original_style_ids: list[int],
) -> None:
    paragraph_snapshot = capture_paragraph_properties(paragraph)

    corrected_style_ids = build_corrected_style_map(
        original_text,
        corrected_text,
        original_style_ids,
    )

    grouped_text = group_text_by_style(
        corrected_text,
        corrected_style_ids,
    )

    paragraph.clear()
    restore_paragraph_properties(paragraph, paragraph_snapshot)

    for text_part, style_id in grouped_text:
        if not text_part:
            continue

        safe_style_id = min(max(style_id, 0), len(styles) - 1)
        run = paragraph.add_run()
        run.text = text_part
        apply_run_style(run, styles[safe_style_id])


# =============================================================================
# 12. TRAITEMENT AVEC PYTHON-PPTX
# =============================================================================

def display_correction(
    slide_number: int,
    element_type: str,
    original_text: str,
    corrected_text: str,
) -> None:
    print()
    print(f"Diapositive {slide_number} — {element_type}")
    print(f"Avant : {original_text}")
    print(f"Après : {corrected_text}")


def process_paragraph(
    paragraph: Any,
    slide_number: int,
    element_type: str,
    language_tool: Any,
    protected_terms: list[str],
    statistics: FileStatistics,
) -> bool:
    (
        original_text,
        styles,
        original_style_ids,
    ) = build_character_style_map(paragraph)

    if not original_text.strip():
        return False

    correction = correct_text(
        original_text,
        language_tool,
        protected_terms,
    )

    statistics.corrections_ignored += correction.ignored_count
    statistics.protected_corrections += correction.protected_count
    statistics.style_corrections_ignored += correction.style_ignored_count
    statistics.corrections_without_suggestion += (
        correction.without_suggestion_count
    )

    if not correction.changed:
        return False

    rebuild_paragraph_with_styles(
        paragraph=paragraph,
        original_text=original_text,
        corrected_text=correction.corrected_text,
        styles=styles,
        original_style_ids=original_style_ids,
    )

    statistics.corrections_applied += correction.applied_count

    display_correction(
        slide_number=slide_number,
        element_type=element_type,
        original_text=original_text,
        corrected_text=correction.corrected_text,
    )

    return True


def process_text_frame(
    text_frame: Any,
    slide_number: int,
    element_type: str,
    language_tool: Any,
    protected_terms: list[str],
    statistics: FileStatistics,
) -> bool:
    changed = False

    for paragraph in text_frame.paragraphs:
        changed = process_paragraph(
            paragraph=paragraph,
            slide_number=slide_number,
            element_type=element_type,
            language_tool=language_tool,
            protected_terms=protected_terms,
            statistics=statistics,
        ) or changed

    return changed


def process_table(
    table: Any,
    slide_number: int,
    language_tool: Any,
    protected_terms: list[str],
    statistics: FileStatistics,
) -> bool:
    changed = False
    processed_cells: set[int] = set()

    for row_index, row in enumerate(table.rows, start=1):
        for column_index, cell in enumerate(row.cells, start=1):
            cell_identity = id(cell._tc)

            # Une cellule fusionnée peut apparaître à plusieurs positions.
            if cell_identity in processed_cells:
                continue

            processed_cells.add(cell_identity)

            changed = process_text_frame(
                text_frame=cell.text_frame,
                slide_number=slide_number,
                element_type=(
                    f"Tableau — ligne {row_index}, colonne {column_index}"
                ),
                language_tool=language_tool,
                protected_terms=protected_terms,
                statistics=statistics,
            ) or changed

    return changed


def process_group_shape(
    group_shape: Any,
    slide_number: int,
    language_tool: Any,
    protected_terms: list[str],
    statistics: FileStatistics,
) -> bool:
    changed = False

    for child_shape in group_shape.shapes:
        changed = process_shape(
            shape=child_shape,
            slide_number=slide_number,
            language_tool=language_tool,
            protected_terms=protected_terms,
            statistics=statistics,
        ) or changed

    return changed


def shape_element_type(shape: Any) -> str:
    if safe_getattr(shape, "is_placeholder", False):
        return "Espace réservé"

    if safe_getattr(shape, "has_table", False):
        return "Tableau"

    return "Zone de texte"


def process_shape(
    shape: Any,
    slide_number: int,
    language_tool: Any,
    protected_terms: list[str],
    statistics: FileStatistics,
) -> bool:
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        return process_group_shape(
            group_shape=shape,
            slide_number=slide_number,
            language_tool=language_tool,
            protected_terms=protected_terms,
            statistics=statistics,
        )

    if safe_getattr(shape, "has_table", False):
        return process_table(
            table=shape.table,
            slide_number=slide_number,
            language_tool=language_tool,
            protected_terms=protected_terms,
            statistics=statistics,
        )

    # Les graphiques sont volontairement ignorés dans la version 1.
    if safe_getattr(shape, "has_chart", False):
        return False

    if safe_getattr(shape, "has_text_frame", False):
        return process_text_frame(
            text_frame=shape.text_frame,
            slide_number=slide_number,
            element_type=shape_element_type(shape),
            language_tool=language_tool,
            protected_terms=protected_terms,
            statistics=statistics,
        )

    return False


def process_slide(
    slide: Any,
    slide_number: int,
    language_tool: Any,
    protected_terms: list[str],
    statistics: FileStatistics,
) -> bool:
    changed = False

    for shape in slide.shapes:
        changed = process_shape(
            shape=shape,
            slide_number=slide_number,
            language_tool=language_tool,
            protected_terms=protected_terms,
            statistics=statistics,
        ) or changed

    statistics.slides_processed += 1
    return changed

def process_presentation(
    presentation: Presentation,
    language_tool: Any,
    protected_terms: list[str],
    temporary_file: Path,
    statistics: FileStatistics,
) -> None:
    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_changed = process_slide(
            slide=slide,
            slide_number=slide_number,
            language_tool=language_tool,
            protected_terms=protected_terms,
            statistics=statistics,
        )

        # Sauvegarde uniquement si cette diapositive a réellement été modifiée.
        if slide_changed:
            presentation.save(str(temporary_file))

# =============================================================================
# 13. VALIDATION ET REMPLACEMENT SÉCURISÉ
# =============================================================================

def validate_corrected_pptx(file_path: Path) -> tuple[bool, str]:
    if not file_path.exists():
        return False, "le fichier temporaire n'existe pas"

    if file_path.stat().st_size == 0:
        return False, "le fichier temporaire est vide"

    is_valid, reason = is_valid_pptx_archive(file_path)

    if not is_valid:
        return False, reason

    try:
        presentation = Presentation(str(file_path))

        if len(presentation.slides) == 0:
            return False, "la présentation ne contient aucune diapositive"

    except Exception as exc:
        return False, (
            "python-pptx ne peut pas rouvrir le fichier corrigé "
            f"({type(exc).__name__}: {exc})"
        )

    return True, ""


def safely_replace_original(
    temporary_file: Path,
    original_file: Path,
) -> None:
    """
    Remplace l'original uniquement après validation.

    os.replace() réalise un remplacement direct lorsque le système et le
    volume de stockage le permettent.
    """
    os.replace(str(temporary_file), str(original_file))


# =============================================================================
# 14. AFFICHAGE DES RÉSULTATS
# =============================================================================

def display_file_summary(statistics: FileStatistics) -> None:
    language_label = {
        "fr-FR": "français",
        "en-US": "anglais",
    }.get(
        statistics.detected_language,
        statistics.detected_language or "non détectée",
    )

    print()
    print("-" * 60)
    print(f"Fichier : {statistics.filename}")
    print(f"Langue : {language_label}")
    print(f"Diapositives traitées : {statistics.slides_processed}")
    print(f"Corrections appliquées : {statistics.corrections_applied}")
    print(f"Corrections ignorées : {statistics.corrections_ignored}")
    print(
        "Corrections sur texte protégé : "
        f"{statistics.protected_corrections}"
    )
    print(
        "Suggestions de style ignorées : "
        f"{statistics.style_corrections_ignored}"
    )
    print(
        "Suggestions sans remplacement : "
        f"{statistics.corrections_without_suggestion}"
    )
    print(f"Erreurs : {statistics.errors}")
    print(f"Statut : {statistics.status}")

    if statistics.skip_reason:
        print(f"Motif : {statistics.skip_reason}")

    print("-" * 60)


def display_final_summary(all_statistics: list[FileStatistics]) -> None:
    corrected = sum(item.status == "CORRECTED" for item in all_statistics)
    unchanged = sum(item.status == "UNCHANGED" for item in all_statistics)
    skipped = sum(item.status == "SKIPPED" for item in all_statistics)
    failed = sum(item.status == "FAILED" for item in all_statistics)

    total_applied = sum(
        item.corrections_applied for item in all_statistics
    )
    total_ignored = sum(
        item.corrections_ignored for item in all_statistics
    )

    print()
    print("=" * 60)
    print("RÉSUMÉ FINAL")
    print("=" * 60)
    print(f"Fichiers trouvés : {len(all_statistics)}")
    print(f"Fichiers corrigés : {corrected}")
    print(f"Fichiers sans correction : {unchanged}")
    print(f"Fichiers ignorés : {skipped}")
    print(f"Fichiers en erreur : {failed}")
    print(f"Corrections appliquées : {total_applied}")
    print(f"Corrections ignorées : {total_ignored}")
    print("=" * 60)


# =============================================================================
# 15. TRAITEMENT D'UN FICHIER
# =============================================================================

def process_powerpoint_file(
    file_path: Path,
    protected_terms: list[str],
    language_tools: dict[str, Any],
) -> FileStatistics:
    statistics = FileStatistics(filename=file_path.name)

    print()
    print("=" * 60)
    print(f"Traitement : {file_path.name}")
    print("=" * 60)

    can_process, reason = can_process_file(file_path)

    if not can_process:
        statistics.status = "SKIPPED"
        statistics.skip_reason = reason
        display_file_summary(statistics)
        return statistics

    try:
        temporary_handle = tempfile.NamedTemporaryFile(
            mode="wb",
            prefix=f".{file_path.stem}.correction_",
            suffix=".pptx",
            dir=file_path.parent,
            delete=False,
        )
        temporary_file = Path(temporary_handle.name)
        temporary_handle.close()

        try:
            shutil.copy2(file_path, temporary_file)

            presentation = Presentation(str(temporary_file))

            detection_text = extract_language_detection_text(presentation)
            detected_language, detection_error = (
                detect_presentation_language(detection_text)
            )

            if detected_language is None:
                statistics.status = "SKIPPED"
                statistics.skip_reason = detection_error
                display_file_summary(statistics)
                return statistics

            statistics.detected_language = detected_language
            language_tool = language_tools[detected_language]

            print(
                "Langue détectée : "
                + (
                    "français"
                    if detected_language == "fr-FR"
                    else "anglais"
                )
            )

            process_presentation(
                presentation=presentation,
                language_tool=language_tool,
                protected_terms=protected_terms,
                temporary_file=temporary_file,
                statistics=statistics,
            )

            # Force la libération de l'objet avant la validation et le
            # remplacement du fichier sous Windows.
            del presentation

            is_valid, validation_error = validate_corrected_pptx(
                temporary_file
            )

            if not is_valid:
                raise RuntimeError(
                    "Échec de la validation du fichier corrigé : "
                    f"{validation_error}"
                )

            if statistics.corrections_applied > 0:
                safely_replace_original(
                    temporary_file=temporary_file,
                    original_file=file_path,
                )
                statistics.status = "CORRECTED"
            else:
                statistics.status = "UNCHANGED"

        finally:
            if temporary_file.exists():
                try:
                    temporary_file.unlink()
                except OSError:
                    pass

    except Exception as exc:
        statistics.errors += 1
        statistics.status = "FAILED"
        statistics.skip_reason = (
            f"{type(exc).__name__}: {exc}"
        )

        print()
        print(f"ERREUR pendant le traitement de {file_path.name}")
        print("Le fichier original n'a pas été modifié.")
        print(
            "Les fichiers temporaires seront supprimés automatiquement."
        )
        print(f"Détail : {type(exc).__name__}: {exc}")

    display_file_summary(statistics)
    return statistics


# =============================================================================
# 16. PROGRAMME PRINCIPAL
# =============================================================================

def main() -> int:
    language_tools: dict[str, Any] = {}

    try:
        validate_configuration()
        protected_terms = load_protected_terms()
        powerpoint_files = find_powerpoint_files()

        if not powerpoint_files:
            print(
                "Aucun fichier .pptx à traiter dans le répertoire :\n"
                f"{POWERPOINT_DIRECTORY}"
            )
            return 0

        language_tools = initialize_language_tools()

        all_statistics: list[FileStatistics] = []

        for file_path in powerpoint_files:
            statistics = process_powerpoint_file(
                file_path=file_path,
                protected_terms=protected_terms,
                language_tools=language_tools,
            )
            all_statistics.append(statistics)

        display_final_summary(all_statistics)
        return 0

    except KeyboardInterrupt:
        print("\nTraitement interrompu par l'utilisateur.")
        return 130

    except Exception as exc:
        print()
        print("ERREUR FATALE")
        print(f"{type(exc).__name__}: {exc}")
        return 1

    finally:
        close_language_tools(language_tools)


if __name__ == "__main__":
    sys.exit(main())
